from data_extraction import *
from text_var import *

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Variables
filename = "/Users/gerarddanjou/Desktop/test_pptx.pptx"


##### CODE HERE #####

# Open

prs = Presentation(filename)

# Data input
slides = prs.slides
first_slide = slides[0]

output_variable, output_variable_name, output_variable_average = compute_average("file_data_3",[3,12,13,14,15,16,17,19])
list_name = []
list_average = []

for i in range (len(output_variable)):
    #txBox = first_slide.shapes.add_textbox(100, 100, 100, 100)
    #tf = txBox.text_frame
    #p = tf.add_paragraph()
    #p.text = output_variable[i]
    
    list_name.append(output_variable_name[i])
    list_average.append(output_variable_average[i])

# define chart data ---------------------
chart_data = CategoryChartData()

chart_data.categories = list_name
chart_data.add_series('DATA :) ', tuple(list_average))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

first_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

# Save and close
prs.save(filename)






