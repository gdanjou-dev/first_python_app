from pptx.util import Inches
from pptx import Presentation

def open_PowerPoint_Presentation(oldFileName, newFileName, img, left, top):
    prs = Presentation(oldFileName)
    slide = prs.slides[3]
    pic = slide.shapes.add_picture(img, left, top)
    prs.save(newFileName)

open_PowerPoint_Presentation('Template.pptx', 'NewTemplate.pptx',
                             'mypic.png', Inches(1), Inches(1))